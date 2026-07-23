import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors (2018 Neutral Enterprise Palette)
    NAVY = RGBColor(26, 37, 47)       # #1a252f
    BLUE = RGBColor(41, 128, 185)     # #2980b9
    SLATE_BG = RGBColor(244, 246, 249) # #f4f6f9
    DARK_TEXT = RGBColor(44, 62, 80)  # #2c3e50
    WHITE = RGBColor(255, 255, 255)
    RED = RGBColor(192, 57, 43)
    GREEN = RGBColor(39, 174, 96)

    blank_slide_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="SECUREBANK PRO TECHNICAL PRESENTATION"):
        # Header background banner
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        
        # Category sub-label
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11.733), Inches(0.3))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = BLUE
        
        # Main Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.6))
        tf2 = tx_box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Title Box
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(2.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "SecureBank Pro"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p_sub = tf.add_paragraph()
    p_sub.text = "Intentionally Vulnerable Banking Portal & Security Education Lab"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = BLUE

    # Links Box
    tb_links = slide1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.2))
    tf_l = tb_links.text_frame
    
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "🌐 Live Deployment: https://securebankpro-vuln-web.vercel.app/"
    p_l1.font.size = Pt(14)
    p_l1.font.bold = True
    p_l1.font.color.rgb = WHITE
    
    p_l2 = tf_l.add_paragraph()
    p_l2.text = "💻 GitHub Repository: https://github.com/Bshashank123/securebankpro-vuln-web"
    p_l2.font.size = Pt(14)
    p_l2.font.bold = True
    p_l2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 2: Project Summary & Purpose
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "Executive Project Summary")
    
    tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p_para = tf2.paragraphs[0]
    p_para.text = "SecureBank Pro is an intentionally vulnerable, full-stack online banking web application engineered for cybersecurity education, vulnerability analysis, and web security defense comparison. Designed with a clean 2018 corporate enterprise aesthetic, the system features a dual-controller architecture that enables users to toggle dynamically between an Attack Mode (exposing six critical OWASP Top 10 vulnerabilities) and a Secure Mode (implementing enterprise-grade defensive mitigations), backed by a real-time telemetry engine that tracks security testing progress and exploit attempts across lab sessions."
    p_para.font.size = Pt(18)
    p_para.font.color.rgb = DARK_TEXT
    p_para.line_spacing = 1.3

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture & Tech Stack
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "System Architecture & Technology Stack")
    
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    
    bullets3 = [
        ("Core Backend:", "Node.js & Express.js server-rendered application architecture."),
        ("Database Engine:", "SQLite3 dual-table schema managing core banking records and telemetry metrics."),
        ("Frontend & UI Design:", "Server-side EJS templates with a 2018 corporate neutral palette (Slate & Navy)."),
        ("Containerization & Cloud:", "Packaged via Docker Compose and deployed to Vercel Serverless Function architecture."),
        ("Session Management:", "Express-Session tracking dynamic Attack Mode vs Secure Mode state per participant.")
    ]
    for i, (title, desc) in enumerate(bullets3):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 4: Interactive Dual-Mode Engine
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "Interactive Dual-Mode Controller Engine")
    
    tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.5))
    tf4 = tb4.text_frame
    tf4.word_wrap = True

    bullets4 = [
        ("Dynamic Mode Switcher:", "Users toggle between ATTACK MODE and SECURE MODE via the top navigation bar."),
        ("Single Endpoint Routing:", "Identical URLs evaluate session state (req.session.isSecureMode) to route requests."),
        ("Attack Mode Branch:", "Executes vulnerable code routines (string concatenation, unescaped HTML, missing CSRF tokens)."),
        ("Secure Mode Branch:", "Executes enterprise defenses (prepared statements, HTML entity encoding, CSRF tokens, lockouts)."),
        ("Side-by-Side Learning:", "Allows developers to immediately test payloads and observe why secure mitigations succeed.")
    ]
    for i, (title, desc) in enumerate(bullets4):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 5: Real-Time Telemetry & Progress Counters
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "Real-Time Telemetry & Progress Tracking Engine")
    
    tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.5))
    tf5 = tb5.text_frame
    tf5.word_wrap = True

    bullets5 = [
        ("Live Home Dashboard:", "Displays total lab sessions, total exploit attempts, and participants who completed all 6 vulns."),
        ("Automated Solve Verification:", "Payloads trigger attempt counts; COMPLETED status is awarded strictly upon successful exploit execution."),
        ("Session Progress Tracking:", "SQLite session_progress table records individual progress per user session."),
        ("Administrative Audit Panel:", "Accessible at /admin for viewing timestamped attack logs and aggregate success rates.")
    ]
    for i, (title, desc) in enumerate(bullets5):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.text = f"• {title} {desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(16)

    # -------------------------------------------------------------
    # SLIDES 6 - 11: Vulnerabilities 1 - 6
    # -------------------------------------------------------------
    vulns_data = [
        ("Vulnerability 1: SQL Injection (SQLi)", "/login (POST)",
         "Concatenates raw user input into SQL query strings (SELECT * FROM users WHERE username = '${username}'...).",
         "Inputting admin' -- bypasses password verification to log in as Administrator.",
         "Uses Parameterized Queries (WHERE username = ? AND password = ?), neutralizing SQL syntax manipulation."),
        
        ("Vulnerability 2: Weak Authentication & Brute-Force", "/login (POST)",
         "Includes default admin account (admin/admin) with zero login throttling or account lockouts.",
         "Enables unlimited automated credential brute-forcing via tools like Burp Suite Intruder.",
         "Enforces session failed login counters, triggering account lockout after 5 consecutive failures."),
        
        ("Vulnerability 3: Stored Cross-Site Scripting (XSS)", "/profile/update (POST) & /profile (GET)",
         "Stores unsanitized bio HTML input into SQLite and renders raw output via EJS <%- user.about_me %>.",
         "Submitting <script>alert('XSS')</script> persistently executes JavaScript whenever profile is loaded.",
         "Enforces EJS contextual entity encoding <%= user.about_me %>, safely escaping script tags into plain text."),
        
        ("Vulnerability 4: Insecure Direct Object Reference (IDOR)", "/transaction?id=... (GET)",
         "Queries database receipts directly by URL parameter id without checking session user ownership.",
         "User john_doe changing ?id=102 to ?id=104 discloses private executive financial receipts.",
         "Validates session ownership (sender_id = user.id OR receiver_id = user.id), returning 403 Forbidden."),
        
        ("Vulnerability 5: Cross-Site Request Forgery (CSRF)", "/transfer (POST)",
         "Processes funds transfers based solely on ambient session cookies without verifying anti-CSRF state tokens.",
         "External malicious websites host auto-submitting forms targeting /transfer to silently drain accounts.",
         "Generates per-session cryptographically signed _csrf tokens and validates presence on POST requests."),
        
        ("Vulnerability 6: Unvalidated File Upload", "/profile/upload (POST)",
         "Accepts arbitrary uploaded file extensions (.txt, .html, .exe) into public directory using original names.",
         "Uploading executable or web script files to serve arbitrary content directly from the server.",
         "Enforces MIME type checking (image/*), restricts extensions (.png, .jpg), and renames files to randomized UUIDs.")
    ]

    for v_title, v_ep, v_impl, v_exp, v_sec in vulns_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        add_header(slide, v_title)
        
        tb_v = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.8))
        tf_v = tb_v.text_frame
        tf_v.word_wrap = True
        
        p0 = tf_v.paragraphs[0]
        p0.text = f"• Target Endpoint: {v_ep}"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = BLUE
        p0.space_after = Pt(10)
        
        p1 = tf_v.add_paragraph()
        p1.text = f"• Attack Mode Implementation: {v_impl}"
        p1.font.size = Pt(15)
        p1.font.color.rgb = RED
        p1.space_after = Pt(10)
        
        p2 = tf_v.add_paragraph()
        p2.text = f"• Exploit Demonstration: {v_exp}"
        p2.font.size = Pt(15)
        p2.font.color.rgb = DARK_TEXT
        p2.space_after = Pt(10)
        
        p3 = tf_v.add_paragraph()
        p3.text = f"• Secure Mode Mitigation: {v_sec}"
        p3.font.size = Pt(15)
        p3.font.color.rgb = GREEN

    # -------------------------------------------------------------
    # SLIDE 12: Deployment & Conclusion (Dark Theme)
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_slide_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = NAVY
    bg12.line.fill.background()

    tb12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Deployment & Educational Conclusion"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)

    concl_points = [
        "Local Execution: Standalone Node.js server (npm start).",
        "Containerized Deployment: Docker Compose for isolated lab environments.",
        "Vercel Serverless Architecture: Deployed with dynamic /tmp SQLite storage.",
        "Educational Value: An interactive hands-on lab demonstrating real-world web application security principles."
    ]
    for pt in concl_points:
        p_c = tf12.add_paragraph()
        p_c.text = f"✓ {pt}"
        p_c.font.size = Pt(18)
        p_c.font.color.rgb = BLUE
        p_c.space_after = Pt(12)

    # Save presentation
    output_filename = "SecureBank_Pro_Presentation.pptx"
    prs.save(output_filename)
    print(f"SUCCESS: Created presentation file '{output_filename}'")

if __name__ == "__main__":
    build_presentation()
